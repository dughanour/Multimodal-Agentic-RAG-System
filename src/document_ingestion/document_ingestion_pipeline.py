from langchain_community.document_loaders import TextLoader, DirectoryLoader, SeleniumURLLoader
from src.embeddings.embedding_service import EmbeddingService
from langchain_core.messages import HumanMessage
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from pathlib import Path
from typing import List, Union, Any
import subprocess
import tempfile
import numpy as np
import pandas as pd
import base64
import pickle
import shutil
import fitz
import io
import os




class DocumentIngestionPipeline:
    """Handles the loading and processing of documents"""

    def __init__(self, chunk_size: int = 900, chunk_overlap: int = 100):
        """
        Initialize the DocumentIngestionPipeline

        Args:
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # Default text splitter (CSV/Excel rows skip chunking)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size = chunk_size,
            chunk_overlap = chunk_overlap,
            separators = ["\n\n", "\n", " ", ""]
        )
    
    def load_from_url(self, url: str)-> List[Document]:
        """
        Load documents from a URL using Playwright

        Args:
            url: URL to load documents from

        Returns:
            List of documents
        """
        loader = SeleniumURLLoader(urls=[url])
        docs = loader.load()
        # normalize metadata for downstream embedding
        for d in docs:
            d.metadata = (d.metadata or {})
            d.metadata.setdefault("type", "text")
            d.metadata.setdefault("source", url)
            d.metadata.setdefault("file_name", url)
        return docs
    
    def load_from_txt(self, file_path: Union[str, Path])-> List[Document]:
        """
        Load documents from a text file

        Args:
            file_path: Path to the text file

        Returns:
            List of documents
        """
        file_path_str = str(file_path)
        loader = TextLoader(file_path_str, encoding="utf-8")
        docs = loader.load()
        for d in docs:
            d.metadata = (d.metadata or {})
            d.metadata.setdefault("type", "text")
            d.metadata.setdefault("source", file_path_str)
            d.metadata.setdefault("file_name", Path(file_path_str).name)
        return docs

    def load_from_txt_dir(self, dir_path: Union[str, Path])-> List[Document]:
        """
        Load documents from a directory of text files

        Args:
            dir_path: Path to the directory of text files

        Returns:
            List of documents
        """
        loader = DirectoryLoader(str(dir_path),loader_cls=TextLoader, loader_kwargs={'encoding': 'utf-8'})
        docs = loader.load()
        for d in docs:
            d.metadata = (d.metadata or {})
            d.metadata.setdefault("type", "text")
            # DirectoryLoader generally sets 'source' to file path
            src = d.metadata.get("source")
            if src:
                d.metadata.setdefault("file_name", Path(src).name)
            else:
                d.metadata.setdefault("source", str(dir_path))
                d.metadata.setdefault("file_name", Path(str(dir_path)).name)
        return docs
    
    def load_from_excel(self, file_path: Union[str, Path], rows_per_doc: int = 3) -> List[Document]:
        """
        Load documents from an Excel file - groups multiple rows per document.
        
        Uses the same format as CSV: "Column1: Value1 | Column2: Value2 | Column3: Value3"
        Rows are separated by newlines within each document.

        Args:
            file_path: Path to the Excel file
            rows_per_doc: Number of complete rows per document (default: 3)

        Returns:
            List of documents (grouped rows, no chunking needed)
        """
        file_path_str = str(file_path)
        documents: List[Document] = []
        
        try:
            # Read Excel with pandas
            df = pd.read_excel(file_path_str)
            
            # Get column names
            columns = df.columns.tolist()
            
            # Format all rows with column headers embedded
            all_row_texts = []
            for idx, row in df.iterrows():
                # Format: "Column: Value | Column: Value | ..."
                row_parts = []
                for col in columns:
                    val = row[col]
                    if pd.notna(val):
                        row_parts.append(f"{col}: {val}")
                
                row_text = " | ".join(row_parts)
                if row_text.strip():
                    all_row_texts.append((idx, row_text))
            
            # Group rows into documents (e.g., 3 rows per document)
            for i in range(0, len(all_row_texts), rows_per_doc):
                group = all_row_texts[i:i + rows_per_doc]
                
                # Combine rows with newlines
                combined_text = "\n".join([text for _, text in group])
                row_indices = [idx for idx, _ in group]
                
                doc = Document(
                    page_content=combined_text,
                    metadata={
                        "type": "excel_row",  # Special type to skip chunking
                        "source": file_path_str,
                        "file_name": Path(file_path_str).name,
                        "row_indices": row_indices,
                        "rows_in_doc": len(group)
                    }
                )
                documents.append(doc)
            
            print(f"✅ Loaded {len(all_row_texts)} rows from Excel as {len(documents)} documents ({rows_per_doc} rows each)")
            return documents
            
        except Exception as e:
            print(f"❌ Error loading Excel {file_path_str}: {e}")
            import traceback
            traceback.print_exc()
            return []
    

    def load_from_csv(self, file_path: Union[str, Path], rows_per_doc: int = 3) -> List[Document]:
        """
        Load documents from a CSV file - groups multiple rows per document.
        
        Each row includes column headers embedded, ensuring the LLM 
        always has context about what each value means.
        
        Format per row: "Column1: Value1 | Column2: Value2 | Column3: Value3"
        Rows are separated by newlines within each document.

        Args:
            file_path: Path to the CSV file
            rows_per_doc: Number of complete rows per document (default: 3)

        Returns:
            List of documents (grouped rows, no chunking needed)
        """
        file_path_str = str(file_path)
        documents: List[Document] = []
        
        try:
            # Read CSV with pandas - handles errors gracefully
            df = pd.read_csv(
                file_path_str,
                on_bad_lines='skip',  # Skip malformed rows
                encoding='utf-8',
                encoding_errors='replace'
            )
            
            # Get column names
            columns = df.columns.tolist()
            
            # Format all rows with column headers
            all_row_texts = []
            for idx, row in df.iterrows():
                # Format: "Column: Value | Column: Value | ..."
                row_parts = []
                for col in columns:
                    val = row[col]
                    if pd.notna(val):
                        row_parts.append(f"{col}: {val}")
                
                row_text = " | ".join(row_parts)
                if row_text.strip():
                    all_row_texts.append((idx, row_text))
            
            # Group rows into documents (e.g., 3 rows per document)
            for i in range(0, len(all_row_texts), rows_per_doc):
                group = all_row_texts[i:i + rows_per_doc]
                
                # Combine rows with newlines
                combined_text = "\n".join([text for _, text in group])
                row_indices = [idx for idx, _ in group]
                
                doc = Document(
                    page_content=combined_text,
                    metadata={
                        "type": "csv_row",  # Special type to skip chunking
                        "source": file_path_str,
                        "file_name": Path(file_path_str).name,
                        "row_indices": row_indices,
                        "rows_in_doc": len(group)
                    }
                )
                documents.append(doc)
            
            print(f"✅ Loaded {len(all_row_texts)} rows from CSV as {len(documents)} documents ({rows_per_doc} rows each)")
            return documents
            
        except Exception as e:
            print(f"❌ Error loading CSV {file_path_str}: {e}")
            import traceback
            traceback.print_exc()
            return []

    def load_img(self, img_path: Union[str, Path]) -> List[Document]:
        """
        Load an image file and wrap it in a LangChain Document, 
        storing its base64-encoded content in metadata.
    
        Args:
            img_path: Path to the image file (.png, .jpg, .jpeg, etc.)
    
        Returns:
            List[Document]: Document representing the image.
        """
        if not os.path.exists(img_path):
            raise FileNotFoundError(f"Image not found: {img_path}")

        # Load the image
        with Image.open(img_path) as img:
            pil_image = img.convert("RGB")
        
        # Convert to base64 for serialization
        buffered = io.BytesIO()
        pil_image.save(buffered, format="PNG")
        img_base64 = base64.b64encode(buffered.getvalue()).decode()

        # Create document
        doc = Document(
            page_content=f"[image:{img_path}]",
            metadata={
                "type": "image",
                "source": str(img_path),
                "file_name": Path(img_path).name,
                "image_base64": img_base64
            }
        )
        return [doc]

        
        
    
    def load_from_pdf(self, file_path: Union[str, Path], strategy: str = "standard") -> List[Document]:
        """
        Load documents from a PDF file with multimodal support (text and images)
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of Document objects containing both text and image documents
        """
        documents = []
        pdf_path = str(file_path)
        
        try:
            # Open PDF with fitz
            doc = fitz.open(pdf_path)
            
            # Process each page
            for page_num, page in enumerate(doc):
                # Extract text from page
                text = page.get_text()
                if text.strip():
                    # Create text document (page_num + 1 for human-readable 1-based page numbers)
                    text_doc = Document(
                        page_content=text,
                        metadata={
                            "page": page_num + 1,
                            "type": "text",
                            "source": pdf_path,
                            "file_name": Path(pdf_path).name
                        }
                    )
                    documents.append(text_doc)
                
                # Extract discrete images ONLY for the 'standard' strategy.
                # The 'summarize' strategy will get visuals from the full-page render.
                if strategy == "standard":
                    for img_index, img in enumerate(page.get_images(full=True)):
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Convert to PIL Image
                            pil_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                            
                            # Create unique identifier
                            image_id = f"page{page_num}_img_{img_index}"
                            
                            # Store image as base64 for later use with LLM
                            buffered = io.BytesIO()
                            pil_image.save(buffered, format="PNG")
                            img_base64 = base64.b64encode(buffered.getvalue()).decode()
                            
                            # Create image document (page_num + 1 for human-readable 1-based page numbers)
                            image_doc = Document(
                                page_content=f"[image:{image_id}]",
                                metadata={
                                    "page": page_num + 1,
                                    "type": "image",
                                    "image_id": image_id,
                                    "source": pdf_path,
                                    "file_name": Path(pdf_path).name,
                                    "image_base64": img_base64,
                                    "is_full_page": False # Flag for discrete images
                                }
                            )
                            documents.append(image_doc)
                            
                        except Exception as e:
                            print(f"Error processing image {img_index} on page {page_num}: {e}")
                            continue
                
                # Render the full page as an image for visual analysis of charts/diagrams
                try:
                    pix = page.get_pixmap(dpi=150)  # Higher DPI for better quality
                    img_bytes = pix.tobytes("png")
                    img_base64 = base64.b64encode(img_bytes).decode("utf-8")

                    full_page_image_doc = Document(
                        page_content="[Full page image for visual analysis]",
                        metadata={
                            "page": page_num + 1,  # +1 for human-readable 1-based page numbers
                            "type": "image",
                            "is_full_page": True,  # Flag for full page renders
                            "source": pdf_path,
                            "file_name": Path(pdf_path).name,
                            "image_base64": img_base64
                        }
                    )
                    documents.append(full_page_image_doc)
                except Exception as e:
                    print(f"Error rendering page {page_num} to image: {e}")
            
            doc.close()
            
        except Exception as e:
            print(f"Error loading PDF {pdf_path}: {e}")
            return []
        
        return documents

    def load_from_pdf_dir(self, dir_path: Union[str, Path], strategy: str = "standard") -> List[Document]:
        """
        Load documents from a directory of PDF files
        
        Args:
            dir_path: Path to directory containing PDF files
            
        Returns:
            List of Document objects from all PDFs in the directory
        """
        documents = []
        dir_path = Path(dir_path)
        
        if not dir_path.exists():
            print(f"Directory {dir_path} does not exist")
            return documents
        
        # Find all PDF files in directory
        pdf_files = list(dir_path.glob("*.pdf"))
        
        if not pdf_files:
            print(f"No PDF files found in {dir_path}")
            return documents
        
        print(f"Found {len(pdf_files)} PDF files in {dir_path}")
        
        for pdf_file in pdf_files:
            print(f"Processing PDF: {pdf_file.name}")
            pdf_docs = self.load_from_pdf(pdf_file, strategy=strategy)
            documents.extend(pdf_docs)
        
        return documents


    
    def _convert_pptx_to_pdf(self, pptx_path: Union[str, Path]) -> str:
        """
        Convert a PPTX file to PDF using LibreOffice.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Path to the generated PDF file, or empty string if conversion failed
        """
        pptx_path = Path(pptx_path).resolve()
        temp_dir = tempfile.mkdtemp(prefix="pptx_to_pdf_")
        
        try:
            # LibreOffice command to convert PPTX to PDF
            cmd = [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                str(pptx_path)
            ]
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120
            )
            
            if result.returncode != 0:
                print(f"LibreOffice conversion failed: {result.stderr}")
                return ""
            
            pdf_name = pptx_path.stem + ".pdf"
            pdf_path = os.path.join(temp_dir, pdf_name)
            
            if os.path.exists(pdf_path):
                return pdf_path
            else:
                print(f"PDF file not found after conversion: {pdf_path}")
                return ""
                
        except FileNotFoundError:
            print("❌ LibreOffice (soffice) not found. Please install LibreOffice.")
            return ""
        except subprocess.TimeoutExpired:
            print("❌ LibreOffice conversion timed out.")
            return ""
        except Exception as e:
            print(f"Error converting PPTX to PDF: {e}")
            return ""


    def load_from_pptx(self, file_path: Union[str, Path], strategy: str = "standard") -> List[Document]:
        """
        Load documents from a PowerPoint (.pptx) file with multimodal support.
        
        For 'summarize' strategy: Converts to PDF and renders full slides as images
        for LLM visual analysis (same as PDF processing).

        Args:
            file_path: Path to the PowerPoint (.pptx) file
            strategy: "standard" extracts discrete images, "summarize" renders full slides

        Returns:
            List[Document] where:
                - text slides become Documents with metadata {"type":"text", "slide_number": ...}
                - each extracted image becomes its own Document with metadata {"type":"image", ...}
        """
        documents: List[Document] = []
        pptx_path = Path(file_path)
        pptx_path_str = str(file_path)

        try:
            prs = Presentation(pptx_path_str)

            for slide_idx, slide in enumerate(prs.slides):
                slide_texts: List[str] = []
                images_docs_local: List[Document] = []
                image_ids_for_slide: List[str] = []
                img_idx = 0

                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                            txt = shape.text.strip()
                            if txt:
                                slide_texts.append(txt)
                    except Exception:
                        continue
                
                # Extract notes (if any)
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame and notes_slide.notes_text_frame.text.strip():
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        slide_texts.append("Notes: " + notes_text)
                except Exception:
                    pass

                # Extract discrete images ONLY for 'standard' strategy
                # The 'summarize' strategy will get visuals from full-slide render
                if strategy == "standard":
                    for shape in slide.shapes:
                        is_picture = getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                        has_image_attr = hasattr(shape, "image")
                        if not (is_picture or has_image_attr):
                            continue

                        try:
                            image = shape.image
                            image_bytes = image.blob

                            pil_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                            image_id = f"{pptx_path.stem}_slide{slide_idx}_img{img_idx}"
                            buffered = io.BytesIO()
                            pil_image.save(buffered, format="PNG")
                            img_base64 = base64.b64encode(buffered.getvalue()).decode()

                            image_doc = Document(
                                page_content=f"[image:{image_id}]",
                                metadata={
                                    "slide_number": slide_idx + 1,  # +1 for human-readable 1-based slide numbers
                                    "type": "image",
                                    "image_id": image_id,
                                    "source": pptx_path_str,
                                    "file_name": pptx_path.name,
                                    "image_base64": img_base64,
                                    "is_full_page": False
                                },
                            )
                            images_docs_local.append(image_doc)
                            image_ids_for_slide.append(image_id)
                            img_idx += 1
                        
                        except Exception as e:
                            print(f"Error extracting image on slide {slide_idx}: {e}")
                            continue

                # Create text document for the slide (+1 for human-readable 1-based slide numbers)
                slide_content = "\n".join(slide_texts).strip()
                if slide_content:
                    text_doc = Document(
                        page_content=slide_content,
                        metadata={
                            "slide_number": slide_idx + 1,
                            "type": "text",
                            "source": pptx_path_str,
                            "file_name": pptx_path.name,
                            "images": image_ids_for_slide,
                        },
                    )
                    documents.append(text_doc)

                # Append discrete image documents (standard mode only)
                documents.extend(images_docs_local)
            
            # === RENDER FULL SLIDES AS IMAGES (ONLY for 'summarize' strategy) ===
            # Convert PPTX to PDF, then render each page as an image for LLM visual analysis
            if strategy == "summarize":
                pdf_path = self._convert_pptx_to_pdf(file_path)
                
                if pdf_path:
                    try:
                        pdf_doc = fitz.open(pdf_path)
                        
                        for slide_idx, page in enumerate(pdf_doc):
                            try:
                                pix = page.get_pixmap(dpi=150)
                                img_bytes = pix.tobytes("png")
                                img_base64 = base64.b64encode(img_bytes).decode("utf-8")
                                
                                full_slide_image_doc = Document(
                                    page_content="[Full slide image for visual analysis]",
                                    metadata={
                                        "slide_number": slide_idx + 1,  # +1 for human-readable 1-based slide numbers
                                        "type": "image",
                                        "is_full_page": True,
                                        "source": pptx_path_str,
                                        "file_name": pptx_path.name,
                                        "image_base64": img_base64
                                    }
                                )
                                documents.append(full_slide_image_doc)
                            except Exception as e:
                                print(f"Error rendering slide {slide_idx} to image: {e}")
                        
                        pdf_doc.close()
                        
                    except Exception as e:
                        print(f"Error processing converted PDF: {e}")
                    finally:
                        # Cleanup temp PDF directory
                        temp_dir = os.path.dirname(pdf_path)
                        shutil.rmtree(temp_dir, ignore_errors=True)
                else:
                    print("⚠️ PPTX to PDF conversion failed. Full slide images not available.")
            
        except Exception as e:
            print(f"Error loading PPTX {pptx_path}: {e}")
            return []
        
        return documents




    def load_documents(self, sources: List[str], strategy: str = "standard") -> List[Document]:
        """
        Load documents from specific sources.
        
        Args:
            sources: List of file paths, directory paths, or URLs to process.
            
        Returns:
            List of loaded documents
        """
        docs: List[Document] = []
        
        for src in sources:
            src_path = Path(src)
            
            # Handle URLs
            if src.startswith("http://") or src.startswith("https://"):
                print(f"Loading from URL: {src}")
                docs.extend(self.load_from_url(src))
            
            # Handle individual files
            elif src_path.is_file():
                print(f"Processing uploaded file: {src}")
                if src_path.suffix.lower() == '.pdf':
                    if strategy == "docling":
                        from src.document_ingestion.docling_strategy import load_from_pdf_docling
                        text_docs_d, image_docs_d = load_from_pdf_docling(src_path)
                        docs.extend(text_docs_d)
                        docs.extend(image_docs_d)
                    else:
                        docs.extend(self.load_from_pdf(src_path, strategy=strategy))
                elif src_path.suffix.lower() in ['.txt', '.text']:
                    docs.extend(self.load_from_txt(src_path))
                elif src_path.suffix.lower() == ".pptx":
                    # For docling strategy, PPTX falls back to summarize (full slide renders)
                    pptx_strategy = "summarize" if strategy == "docling" else strategy
                    docs.extend(self.load_from_pptx(src_path, strategy=pptx_strategy))
                elif src_path.suffix.lower() in [".xlsx", ".xls"]:
                    docs.extend(self.load_from_excel(src_path))
                elif src_path.suffix.lower() == ".csv":
                    docs.extend(self.load_from_csv(src_path))
                elif src_path.suffix.lower() == ".png" or src_path.suffix.lower() == ".jpg" or src_path.suffix.lower() == ".jpeg":
                    docs.extend(self.load_img(src_path))
                else:
                    print(f"Unsupported file type: {src_path.suffix}")
            
            # Handle directories (when user uploads a folder)
            elif src_path.is_dir():
                print(f"Processing uploaded directory: {src}")
                
                # Load PDFs from directory
                pdf_docs = self.load_from_pdf_dir(src_path, strategy=strategy)
                docs.extend(pdf_docs)
                
                # Load text files from directory
                txt_docs = self.load_from_txt_dir(src_path)
                docs.extend(txt_docs)
            
            else:
                print(f"Source not found: {src}")
        
        print(f"Total documents processed: {len(docs)}")
        return docs


    def _get_image_prompt(self, is_full_page: bool) -> str:
        """Returns the appropriate prompt for image summarization."""
        if is_full_page:
            # For full page images using summarize strategy
            return (
                "Analyze this document page image. Follow these rules strictly:\n\n"
                "1. If the page contains ONLY a title, logo, landing page, section divider, or decorative content "
                "with no charts, graphs, diagrams, or data tables, respond with EXACTLY: "
                "[No significant visual elements to analyze]\n\n"
                "2. If the page contains charts, graphs, diagrams, or data tables, write a brief summary "
                "(maximum 80 words) covering ONLY the key data points, trends, and conclusions. "
                "Mention specific numbers and labels visible in the visuals. "
                "Do NOT describe logos, icons, branding, or decorative elements. "
                "Do NOT repeat text that is already readable on the page."
            )
        else:
            # For discrete images using docling strategy
            return(
                "Analyze and Briefly describe this image in 80 words or less.. Follow these rules strictly:\n\n"
                "1. If the image contains charts, graphs, diagrams, or data tables, write a brief summary "
                "2. (maximum 50-60 words) covering ONLY and ALL the key data points, trends, and conclusions."
                "3. Mention specific numbers and labels visible in the visuals if exists. "
                "Do NOT describe logos, icons, branding, or decorative elements. "

            )

    def _generate_image_summary(self, doc: Document, llm: Any) -> str:
        """
        Generates a textual summary for an image document using a multimodal LLM.
        """
        if not hasattr(doc, "metadata") or "image_base64" not in doc.metadata:
            return "[No image data available]"

        img_base64 = doc.metadata["image_base64"]
        is_full_page = doc.metadata.get("is_full_page", False)

        prompt = self._get_image_prompt(is_full_page)

        multimodal_message = HumanMessage(content=[
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_base64}"}}
        ])

        try:
            response = llm.invoke([multimodal_message])
            return response.content
        except Exception as e:
            print(f"Error generating image summary: {e}")
            return "[Image summary generation failed]"

                
        

    def embed_documents_with_summaries(self, documents: List[Document], embedding_service: EmbeddingService, llm: Any) -> tuple[List[Document], np.ndarray]:
        """
        Compute embeddings for a mixed list of text and image Documents.
        - Text docs: use page_content with embed_text
        - Image docs: generate summaries with llm.batch() (parallel), then embed the summary text.
        """
        embeddings: List[np.ndarray] = []
        docs_to_embed: List[Document] = []

        # Separate image and text documents
        image_docs = [d for d in documents if (d.metadata or {}).get("type") == "image"]
        text_docs = [d for d in documents if (d.metadata or {}).get("type") != "image"]

        # --- Batch summarize all images using llm.batch() (parallel LLM calls) ---
        if image_docs:
            # Build all messages for the batch
            messages_list = []
            valid_image_docs = []
            for doc in image_docs:
                if not hasattr(doc, "metadata") or "image_base64" not in doc.metadata:
                    continue
                img_base64 = doc.metadata["image_base64"]
                is_full_page = doc.metadata.get("is_full_page", False)
                prompt = self._get_image_prompt(is_full_page)

                multimodal_message = HumanMessage(content=[
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_base64}"}}
                ])
                messages_list.append([multimodal_message])
                valid_image_docs.append(doc)

            # Determine if LLM is Groq (use sequential invoke) or Ollama (use batch)
            from src.services.llm_factory import llm_factory
            use_sequential = llm_factory.get_config().get("provider", "").lower() == "groq"

            responses = []
            if use_sequential:
                # Groq: invoke one by one
                for idx, messages in enumerate(messages_list):
                    print(f"🖼️ Summarizing image {idx + 1} of {len(messages_list)} using llm.invoke()...")
                    try:
                        response = llm.invoke(messages)
                        responses.append(response)
                    except Exception as e:
                        print(f"Error summarizing image {idx + 1}: {e}")
                        responses.append(None)
            else:
                # Ollama (or other): use llm.batch() in chunks of 6 images at a time
                BATCH_SIZE = 6
                for i in range(0, len(messages_list), BATCH_SIZE):
                    batch_chunk = messages_list[i : i + BATCH_SIZE]
                    print(f"🖼️ Summarizing images {i + 1}-{i + len(batch_chunk)} of {len(messages_list)} using llm.batch()...")
                    try:
                        batch_responses = llm.batch(batch_chunk)
                        responses.extend(batch_responses)
                    except Exception as e:
                        print(f"Error in batch image summarization: {e}")
                        responses.extend([None] * len(batch_chunk))

            # Process each summary
            for doc, response in zip(valid_image_docs, responses):
                if response is None:
                    summary = "[Image summary generation failed]"
                else:
                    summary = response.content

                # CRITICAL: If the summary indicates no visuals, skip this document entirely.
                if summary.strip() == "[No significant visual elements to analyze]":
                    page_num = doc.metadata.get('page', 'N/A')
                    print(f"📄 Skipping full-page image analysis for page {page_num} (no significant visuals found).")
                    continue

                # Keep image summary as a single document (don't chunk)
                # Each image represents one visual element and should be retrievable as a whole
                summary_doc = Document(page_content=summary, metadata=doc.metadata)
                emb = embedding_service.embed_text(summary, task_type="search_document")
                embeddings.append(emb)
                docs_to_embed.append(summary_doc)

        # --- Embed text documents (unchanged) ---
        for doc in text_docs:
            text_content = doc.page_content or ""
            if not text_content.strip():
                continue # skip empty text
            
            emb = embedding_service.embed_text(text_content, task_type="search_document")
            embeddings.append(emb)
            docs_to_embed.append(doc)
                
        return docs_to_embed, np.array(embeddings)


    def process_and_embed_with_summaries(self, sources: List[str], embedding_service: EmbeddingService, llm: Any, strategy: str = "standard", known_sources: set = None) -> tuple[List[Document], np.ndarray]:
        """
        Load sources, generate summaries for images, split text, and compute embeddings.
        """
        docs = self.load_documents(sources, strategy=strategy)
        for doc in docs:
            doc.metadata["ingestion_strategy"] = strategy
        
        # Separate by type: text, csv_row, excel_row, and image
        text_docs = [d for d in docs if (d.metadata or {}).get("type") == "text"]
        csv_row_docs = [d for d in docs if (d.metadata or {}).get("type") == "csv_row"]
        excel_row_docs = [d for d in docs if (d.metadata or {}).get("type") == "excel_row"]
        image_docs = [d for d in docs if (d.metadata or {}).get("type") == "image"]
        
        # Split text documents (CSV/Excel rows skip splitting via split_documents logic)
        # For docling strategy, text is already chunked by Docling's HybridChunker
        all_text_docs = text_docs + csv_row_docs + excel_row_docs
        if strategy == "docling":
            split_text_docs = all_text_docs
        else:
            split_text_docs = self.split_documents(all_text_docs) if all_text_docs else []
        
        # All docs that will be passed to the embedding function
        combined_docs: List[Document] = split_text_docs + image_docs

        # --- Check for already embedded docs --- //For FAISS//
        # Old: FAISS pickle-based deduplication
        # known_sources = set()
        # docs_path = os.path.join(persist_directory, "docs.pkl")
        # if os.path.exists(docs_path):
        #     try:
        #         with open(docs_path,"rb") as f:
        #             existing_docs = pickle.load(f)
        #             for d in existing_docs:
        #                 known_sources.add(d.metadata.get("source"))
        #     except Exception as e:
        #         print(f"⚠️ Could not load existing docs: {e}")

        # New: known_sources is passed by the caller (from PostgreSQL via get_existing_sources())
        if known_sources is None:
            known_sources = set()
        
        new_docs = [d for d in combined_docs if d.metadata.get("source") not in known_sources]

        if not new_docs:
            print("✅ All uploaded documents already exist in vector store. Skipping re-embedding.")
            return [], np.array([])
        
        print(f"🔍 Found {len(new_docs)} new documents to process and embed with summaries.")

        # This new function handles the logic of summarizing images before embedding
        final_docs, embeddings = self.embed_documents_with_summaries(new_docs, embedding_service, llm)
        return final_docs, embeddings


    def split_documents(self, documents: List[Document]) -> List[Document]:
        """
        Split documents into chunks.
        - CSV/Excel rows: NO splitting (each row group is self-contained with headers)
        - Other files: Normal chunks
        
        Args:
            documents: List of documents to split
            
        Returns:
            List of split documents
        """
        row_docs = []  # CSV and Excel rows skip chunking entirely
        other_docs = []
        
        for d in documents:
            doc_type = d.metadata.get("type", "")
            
            # CSV and Excel rows are already complete - no splitting needed
            if doc_type in ("csv_row", "excel_row"):
                row_docs.append(d)
            else:
                other_docs.append(d)
        
        # Split only non-row documents
        split_other = self.text_splitter.split_documents(other_docs) if other_docs else []
        
        # Row docs are returned as-is (no chunking)
        return split_other + row_docs


    def embed_documents(self, documents: List[Document], embedding_service: EmbeddingService) -> np.ndarray:
        """
        Compute embeddings for a mixed list of text and image Documents using CLIP OR NOMIC.
        - Text docs: use page_content with embed_text
        - Image docs: use metadata["image_base64"] with embed_image

        Args:
            documents: List of documents to embed
            embedding_service: Embedding service to use

        Returns:
            List of embeddings
        """
        embeddings: List[np.ndarray] = []
        for doc in documents:
            doc_type = (doc.metadata or {}).get("type", "text")
            if doc_type == "image":
                img_b64 = (doc.metadata or {}).get("image_base64")
                if not img_b64:
                    # skip images without data
                    continue
                try:
                    image_bytes = base64.b64decode(img_b64)
                    pil_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    emb = embedding_service.embed_image(pil_image)
                except Exception:
                    # skip problematic images
                    continue
            else:
                text_content = doc.page_content or ""
                if not text_content.strip():
                    # skip empty text
                    continue
                emb = embedding_service.embed_text(text_content)
            embeddings.append(emb)
        return np.array(embeddings)


    def process_and_embed(self, sources: List[str], embedding_service: EmbeddingService, strategy: str = "standard", known_sources: set = None) -> tuple[List[Document], np.ndarray]:
        """
        Load sources, split only text documents, keep images, then compute embeddings.

        Args:
            sources: List of sources to load documents from
            embedding_service: Embedding service to use
            persist_directory: Path to the directory to persist the documents and embeddings

        Returns:
            Tuple of (documents_in_order, embeddings_array)
        """
        docs = self.load_documents(sources, strategy=strategy)
        for doc in docs:
            doc.metadata["ingestion_strategy"] = strategy

        # Separate by type: text, csv_row, excel_row, and image
        text_docs = [d for d in docs if (d.metadata or {}).get("type") == "text"]
        csv_row_docs = [d for d in docs if (d.metadata or {}).get("type") == "csv_row"]
        excel_row_docs = [d for d in docs if (d.metadata or {}).get("type") == "excel_row"]
        image_docs = [d for d in docs if (d.metadata or {}).get("type") == "image"]
        
        # Split text documents (CSV/Excel rows skip splitting via split_documents logic)
        all_text_docs = text_docs + csv_row_docs + excel_row_docs
        text_docs = self.split_documents(all_text_docs) if all_text_docs else []
        combined_docs: List[Document] = text_docs + image_docs


        # === Check for already embedded docs === //For FAISS//
        # Old: FAISS pickle-based deduplication
        # known_sources = set()
        # docs_path = os.path.join(persist_directory, "docs.pkl")
        # if os.path.exists(docs_path):
        #     try:
        #         with open(docs_path,"rb") as f:
        #             existing_docs = pickle.load(f)
        #             for d in existing_docs:
        #                 known_sources.add(d.metadata.get("source"))
        #     except Exception as e:
        #         print(f"⚠️ Could not load existing docs: {e}")

        # New: known_sources is passed by the caller (from PostgreSQL via get_existing_sources())
        if known_sources is None:
            known_sources = set()
        
        new_docs = [d for d in combined_docs if d.metadata.get("source") not in known_sources]

        if not new_docs:
            print("✅ All uploaded documents already exist in vector store. Skipping re-embedding.")
            return [], np.array([])
        
        print(f"🔍 Found {len(new_docs)} new documents to embed.")


        embeddings = self.embed_documents(new_docs, embedding_service)
        return new_docs, embeddings



#######################################################################################


    
    def process_documents(self, urls: List[str]) -> List[Document]:
        """
        Complete pipeline to load and split documents
        
        Args:
            urls: List of URLs to process
            
        Returns:
            List of processed document chunks
        """
        docs = self.load_documents(urls)
        return self.split_documents(docs)
